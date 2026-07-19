import os
try:
    import openpyxl
except ImportError:
    pass

try:
    from docx import Document
except ImportError:
    pass

try:
    from pptx import Presentation
except ImportError:
    pass

def extract_docx(filepath):
    print(f"--- Extracting {filepath} ---")
    try:
        doc = Document(filepath)
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                print(para.text)
    except Exception as e:
        print(f"Error reading docx: {e}")
    print("-" * 40)

def extract_pptx(filepath):
    print(f"--- Extracting {filepath} ---")
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except Exception as e:
        print(f"Error reading pptx: {e}")
    print("-" * 40)

def extract_excel_info(filepath):
    print(f"--- Extracting Excel {filepath} ---")
    try:
        wb = openpyxl.load_workbook(filepath, data_only=False)
        print("Sheet names:", wb.sheetnames)
        
        # We need to understand the inputs and formulas
        for sheet_name in wb.sheetnames:
            print(f"\nSheet: {sheet_name}")
            ws = wb[sheet_name]
            # Just print the first 30 rows and 10 columns that contain data
            for row in ws.iter_rows(min_row=1, max_row=30, min_col=1, max_col=10):
                row_data = []
                for cell in row:
                    if cell.value is not None:
                        row_data.append(f"{cell.coordinate}: {cell.value}")
                if row_data:
                    print(" | ".join(row_data))
    except Exception as e:
        print(f"Error reading excel: {e}")
    print("-" * 40)

if __name__ == "__main__":
    extract_docx("Nivra Template.docx")
    extract_pptx("Nivra Website Requirements.pptx")
    extract_excel_info("Nivra Goal - Compute SIP_or_StepUP_SIP v3.xlsm")
